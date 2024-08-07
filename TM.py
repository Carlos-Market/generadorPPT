from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from werkzeug.utils import secure_filename
from PIL import Image
import os
from datetime import datetime
from io import BytesIO

app = Flask(__name__)

# Configuración de carpetas y extensiones permitidas
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def formato_miles(numero):
    return '{:,}'.format(numero).replace(',', '.')

@app.route('/')
def index():
    ubicaciones = [
        "Alcalá", "Avenida 68", "Banderas", "Calle 100", "Calle 127", "Calle 146",
        "Calle 45", "Calle 76", "Calle 85", "El Campín", "Héroes", "Movistar arena",
        "NQS Calle 38a Sur", "Paloquemao", "Pepe sierra", "Portal 80", "Portal Américas",
        "Portal El Dorado", "Portal Norte", "Portal Suba", "Portal Sur", "Suba Calle 100",
        "Toberin", "Universidades", "Virrey"
    ]

    options = ''.join([f'<option value="{ubicacion}">{ubicacion}</option>' for ubicacion in ubicaciones])

    return '''
    <!doctype html>
    <html lang="es">
    <head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1, shrink-to-fit=no">
    <title>REPORTES</title>
    <link rel="stylesheet" href="https://stackpath.bootstrapcdn.com/bootstrap/4.5.2/css/bootstrap.min.css">
    </head>
    <body>
    <div class="container">
        <div class="row justify-content-center">
            <div class="col-md-8">
                <h1 class="mt-5 mb-4 text-center">GENERA TU REPORTE DE CAMPAÑA</h1>
                <form id="report-form" method="post" enctype="multipart/form-data" action="/upload">
                    <div class="form-group">
                        <select class="form-control" id="Mes" name="Mes" required>
                            <option value="">Selecciona mes</option>
                            <option value="Enero">Enero</option>
                            <option value="Febrero">Febrero</option>
                            <option value="Marzo">Marzo</option>
                            <option value="Abril">Abril</option>
                            <option value="Mayo">Mayo</option>
                            <option value="Junio">Junio</option>
                            <option value="Julio">Julio</option>
                            <option value="Agosto">Agosto</option>
                            <option value="Septiembre">Septiembre</option>
                            <option value="Octubre">Octubre</option>
                            <option value="Noviembre">Noviembre</option>
                            <option value="Diciembre">Diciembre</option>
                        </select>
                    </div>
                    <div id="ubicaciones-container">
                        <div class="ubicacion-group">
                            <div class="form-group">
                                <select class="form-control" name="Ubicacion[]" required onchange="updateElemento(this)">
                                    <option value="">Selecciona ubicación</option>
                                    ''' + options + '''
                                </select>
                            </div>
                            <div class="form-group">
                                <input type="text" class="form-control" name="Elemento[]" id="elemento" placeholder="Tipo de elemento" readonly>
                            </div>
                            
                            <div class="form-group">
                                <input type="text" class="form-control" name="Trafico[]" placeholder="Trafico" required>
                            </div>

                            <div class="form-group">
                                <label>Sube tus imágenes</label>
                                <input type="file" class="form-control-file" name="files_0[]" multiple required>
                            </div>
                        </div>
                    </div>
                    <button type="button" class="btn btn-secondary" onclick="addUbicacion()">Añadir otra ubicación</button>
                    <button type="submit" class="btn btn-primary btn-block">Crea tu presentación</button>
                </form>
            </div>
        </div>
    </div>
    <script>
        let ubicacionCounter = 1;
        function addUbicacion() {
            var container = document.getElementById('ubicaciones-container');
            var newGroup = document.createElement('div');
            newGroup.classList.add('ubicacion-group');
            newGroup.innerHTML = `
                <div class="form-group">
                    <select class="form-control" name="Ubicacion[]" required onchange="updateElemento(this)">
                        <option value="">Selecciona ubicación</option>
                        ''' + options + '''
                    </select>
                </div>
                <div class="form-group">
                    <input type="text" class="form-control" name="Elemento[]" id="elemento" placeholder="Tipo de elemento" readonly>
                </div>
                
                <div class="form-group">
                    <input type="text" class="form-control" name="Trafico[]" placeholder="Trafico" required>
                </div>

                <div class="form-group">
                    <label>Sube tus imágenes</label>
                    <input type="file" class="form-control-file" name="files_${ubicacionCounter}[]" multiple required>
                </div>`;
            container.appendChild(newGroup);
            ubicacionCounter++;
        }

        function updateElemento(select) {
            var elementoInput = select.parentElement.nextElementSibling.querySelector('input');
            if (select.value === 'Alcalá' || select.value === 'Avenida 68' || select.value === 'Banderas' ||
                select.value === 'Calle 127' || select.value === 'Calle 146' || select.value === 'Calle 45' ||
                select.value === 'Calle 76' || select.value === 'Calle 85' || select.value === 'El Campín' ||
                select.value === 'Movistar arena' || select.value === 'NQS Calle 38a Sur' || select.value === 'Paloquemao' ||
                select.value === 'Pepe sierra' || select.value === 'Suba Calle 100' || select.value === 'Toberin' ||
                select.value === 'Virrey') {
                elementoInput.value = 'Tótem digital';
            } else if (select.value === 'Calle 100' || select.value === 'Héroes' || select.value === 'Portal 80' ||
                    select.value === 'Portal Américas' || select.value === 'Portal El Dorado' || select.value === 'Portal Norte' ||
                    select.value === 'Portal Suba' || select.value === 'Portal Sur' || select.value === 'Universidades') {
                elementoInput.value = 'Pantalla Gran Formato';
            } else {
                elementoInput.value = '';
            }
        }
    </script>
    </body>
    </html>
    '''

@app.route('/upload', methods=['POST'])
def upload_files():
    mes = request.form.get('Mes')
    mes = mes.upper()
    
    ubicaciones = request.form.getlist('Ubicacion[]')
    elementos = request.form.getlist('Elemento[]')
    
    trafico = request.form.getlist('Trafico[]')
    
    # Formatear 'trafico' con separadores de miles
    # trafico = formato_miles(int(trafico))
    
    # Crear presentación PowerPoint
    ppt = Presentation()

    # Configurar el tamaño de la presentación
    ppt.slide_width = Inches(13.334646)
    ppt.slide_height = Inches(7.5)

    # Rutas de las imágenes específicas
    pagina1_path = os.path.join(os.path.dirname(__file__), 'img/1.png')
    pagina2_path = os.path.join(os.path.dirname(__file__), 'img/2.png')
    pagina3_path = os.path.join(os.path.dirname(__file__), 'img/3.png')
    penultima_path = os.path.join(os.path.dirname(__file__), 'img/penultima.png')
    ultima_pagina_path = os.path.join(os.path.dirname(__file__), 'img/ultima.png')

    # Crear las primeras tres diapositivas fijas
    slide_layout = ppt.slide_layouts[6]  # Layout en blanco
    slide = ppt.slides.add_slide(slide_layout)
    imagen_presentacion(slide, pagina1_path)
    
    slide = ppt.slides.add_slide(slide_layout)
    imagen_presentacion(slide, pagina2_path)
    
    # Procesar todas las ubicaciones y sus archivos
    for i, (ubicacion, elemento, trafico) in enumerate(zip(ubicaciones, elementos, trafico)):
        file_key = f'files_{i}[]'
        files = request.files.getlist(file_key)
        
        if not files or not all(allowed_file(file.filename) for file in files):
            return 'No files uploaded or invalid file type', 400
        
        for file in files:
            if allowed_file(file.filename):
                filename = secure_filename(file.filename)
                file_bytes = file.read()  # Leer el archivo en memoria
                img_stream = BytesIO(file_bytes)
                
                try:
                    img = Image.open(img_stream)  # Validar la imagen con PIL
                    img.verify()  # Esto asegura que la imagen está completa y no está corrupta

                    img_stream.seek(0)  # Volver al inicio del archivo en memoria para poder utilizarlo
                    slide = ppt.slides.add_slide(slide_layout)  # Agregar diapositiva en blanco
                    imagen_presentacion(slide, pagina3_path)
                    slide.shapes.add_picture(img_stream, top=Inches(1.2), left=Inches(4), height=Inches(5.1), width=Inches(9))
                    info_foto(slide, elemento, ubicacion, mes, trafico)  # Nota el cambio en el orden aquí
                except (IOError, SyntaxError) as e:
                    return f'Error al procesar la imagen: {filename}', 400
                
    # Añadir la antepenúltima diapositiva fija
    slide = ppt.slides.add_slide(slide_layout)
    imagen_presentacion(slide, penultima_path)

    # Añadir la última diapositiva fija
    slide = ppt.slides.add_slide(slide_layout)
    imagen_presentacion(slide, ultima_pagina_path)

    # Guardar la presentación en memoria
    ppt_io = BytesIO()
    ppt.save(ppt_io)
    ppt_io.seek(0)

    # Proporcionar el archivo para descarga
    return send_file(ppt_io, as_attachment=True, download_name='Reporte.pptx', mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

def imagen_presentacion(slide, img_path):
    img_left = Inches(0)
    img_top = Inches(0)
    img_width = Inches(13.334646)
    img_height = Inches(7.5)
    slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)

def info_foto(slide, elemento, ubicacion, mes, trafico):
    textbox = slide.shapes.add_textbox(top=Inches(2.12), left=Inches(0.5), height=Inches(3.5), width=Inches(1))
    text_frame = textbox.text_frame

    p1 = text_frame.add_paragraph()
    p1.text = "ELEMENTO "
    p1.font.bold = True
    p1.font.size = Pt(16)
    p1.font.color.rgb = RGBColor(153, 146, 144)
    
    p2 = text_frame.add_paragraph()
    p2.text = elemento
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0, 0, 0) 
    
     # Párrafo vacío para separación
    text_frame.add_paragraph().text = ""

    p3 = text_frame.add_paragraph()
    p3.text = "UBICACIÓN "
    p3.font.bold = True
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(153, 146, 144)
    
    p4 = text_frame.add_paragraph()
    p4.text = ubicacion
    p4.font.size = Pt(16)
    p4.font.color.rgb = RGBColor(0, 0, 0) 
    
     # Párrafo vacío para separación
    text_frame.add_paragraph().text = ""
    
    p5 = text_frame.add_paragraph()
    p5.text = "TRAFICO A " + mes
    p5.font.bold = True
    p5.font.size = Pt(16)
    p5.font.color.rgb = RGBColor(153, 146, 144)
    
    p6 = text_frame.add_paragraph()
    p6.text = trafico
    p6.font.size = Pt(16)
    p6.font.color.rgb = RGBColor(0, 0, 0) 
    
    # Cuadro de texto para la fecha
    fecha_actual = datetime.now().strftime("%d/%m/%Y")
    fecha_textbox = slide.shapes.add_textbox(top=Inches(5.65), left=Inches(11.72), width=Inches(1.2), height=Inches(0.2))
    fecha_text_frame = fecha_textbox.text_frame

    p_fecha = fecha_text_frame.add_paragraph()
    p_fecha.text = fecha_actual
    p_fecha.font.size = Pt(16)
    p_fecha.font.color.rgb = RGBColor(255, 255, 255)

if __name__ == '__main__':
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    app.run(debug=True)
