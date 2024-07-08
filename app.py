from flask import Flask, request, send_file, redirect, url_for
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from werkzeug.utils import secure_filename
from io import BytesIO

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    return '''
    <!doctype html>
<html lang="es">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1, shrink-to-fit=no">
    <title>REPORTES</title>
    <link rel="stylesheet" href="https://stackpath.bootstrapcdn.com/bootstrap/4.5.2/css/bootstrap.min.css">
</head>
<body>
    <div class="container">
        <div class="row justify-content-center">
            <div class="col-md-8">
                <h1 class="mt-5 mb-4 text-center">GENERA TU REPORTE DE CAMPAÑA</h1>
                <form method="post" enctype="multipart/form-data" action="/upload">
                    <div class="form-group">
                        <input type="text" class="form-control" id="Cliente" name="Cliente" placeholder="Nombre Cliente" required>
                    </div>
                    <div class="form-group">
                        <select class="form-control" id="Mes" name="Mes" required>
                            <option value="">Selecciona un mes</option>
                            <option value="Enero">Enero</option>
                            <option value="Febrero">Febrero</option>
                            <option value="Marzo">Marzo</option>
                            <option value="Abril">Abril</option>
                            <option value="Mayo">Mayo</option>
                            <option value="Junio">Junio</option>
                            <option value="Julio">Julio</option>
                            <option value="Agosto">Agosto</option>
                            <option value="Septiembre">Septiembre</option>
                            <option value="Octubre">Octubre</option>
                            <option value="Noviembre">Noviembre</option>
                            <option value="Diciembre">Diciembre</option>
                        </select>
                    </div>
                    <div class="form-group">
                        <input type="text" class="form-control" id="Pantalla" name="Pantalla" placeholder="Nombre Pantalla" required>
                    </div>
                    <div class="form-group">
                        <input type="text" class="form-control" id="Elemento" name="Elemento" placeholder="Tipo elemento" required>
                    </div>
                    <div class="form-group">
                        <label for="files">Sube tus imágenes</label>
                        <input type="file" class="form-control-file" id="files" name="files" multiple required>
                    </div>
                    <button type="submit" class="btn btn-primary btn-block">Crea tu presentación</button>
                </form>
            </div>
        </div>
    </div>
    <script src="https://code.jquery.com/jquery-3.5.1.slim.min.js"></script>
    <script src="https://cdn.jsdelivr.net/npm/@popperjs/core@2.5.4/dist/umd/popper.min.js"></script>
    <script src="https://stackpath.bootstrapcdn.com/bootstrap/4.5.2/js/bootstrap.min.js"></script>
    '''

@app.route('/upload', methods=['POST'])
def upload_files():
    cliente = request.form.get('Cliente')
    cliente = cliente.upper()
    
    mes = request.form.get('Mes')
    mes = mes.upper()
    
    pantalla = request.form.get('Pantalla')
    pantalla = pantalla.upper()
    
    elemento = request.form.get('Elemento')
    elemento = elemento.upper()
    files = request.files.getlist('files')
    
    if not files:
        return 'No files uploaded', 400

    # Crear presentación PowerPoint
    ppt = Presentation()

    # Configurar el tamaño de la presentación
    ppt.slide_width = Inches(13.334646)
    ppt.slide_height = Inches(7.5)

    # Rutas de las imágenes específicas
    pagina1_path = os.path.join(os.path.dirname(__file__), 'img/1.png')
    pagina2_path = os.path.join(os.path.dirname(__file__), 'img/2.png')
    pagina3_path = os.path.join(os.path.dirname(__file__), 'img/3.png')
    ultima_pagina_path = os.path.join(os.path.dirname(__file__), 'img/ultima.png')
    marca_agua_path = os.path.join(os.path.dirname(__file__), 'img/marca_agua.png')

    # Crear las primeras tres diapositivas fijas
    slide_layout = ppt.slide_layouts[6]  # Layout en blanco
    slide = ppt.slides.add_slide(slide_layout)
    imagen_presentacion(slide, pagina1_path)
    
    slide = ppt.slides.add_slide(slide_layout)
    imagen_presentacion(slide, pagina2_path)
    datos_reporte(slide, cliente, mes)
    
    slide = ppt.slides.add_slide(slide_layout)
    imagen_presentacion(slide, pagina3_path)
    nombre_comercial(slide, pantalla)

    # Añadir las imágenes subidas a partir de la cuarta diapositiva
    for index, file in enumerate(files):
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            slide = ppt.slides.add_slide(slide_layout)  # Agregar diapositiva en blanco
            pic = slide.shapes.add_picture(filepath, Inches(1), Inches(1), height=Inches(6.29921), width=Inches(11.3))
            agregar_marca_agua(slide, marca_agua_path)
            info_foto(slide, pantalla, elemento)

    # Añadir la última diapositiva fija
    slide = ppt.slides.add_slide(slide_layout)
    imagen_presentacion(slide, ultima_pagina_path)

    # Guardar la presentación en memoria
    ppt_io = BytesIO()
    ppt.save(ppt_io)
    ppt_io.seek(0)

    # Proporcionar el archivo para descarga
    return send_file(ppt_io, as_attachment=True, download_name='Reporte.pptx', mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

def agregar_marca_agua(slide, marca_agua_path):
    marca_agua_left = Inches(12.20)
    marca_agua_top = Inches(0.20)
    marca_agua_width = Inches(0.75)
    marca_agua_height = Inches(0.75)
    slide.shapes.add_picture(marca_agua_path, marca_agua_left, marca_agua_top, width=marca_agua_width, height=marca_agua_height)

def imagen_presentacion(slide, img_path):
    img_left = Inches(0)
    img_top = Inches(0)
    img_width = Inches(13.334646)
    img_height = Inches(7.5)
    slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)

def datos_reporte(slide, cliente, mes):
    textbox = slide.shapes.add_textbox(Inches(3.2), Inches(2.4), Inches(8), Inches(3))
    text_frame = textbox.text_frame

    p1 = text_frame.add_paragraph()
    p1_run = p1.add_run()
    p1_run.text = "REPORTE FOTOGRAFICO"
    p1.font.bold = True
    p1.font.size = Pt(43)

    p2 = text_frame.add_paragraph()
    p2_run = p2.add_run()
    p2_run.text = "CLIENTE: " + cliente
    p2.font.bold = True
    p2.font.size = Pt(43)

    p3 = text_frame.add_paragraph()
    p3_run = p3.add_run()
    p3_run.text = "MES: " + mes
    p3.font.bold = True
    p3.font.size = Pt(43)

def nombre_comercial(slide , pantalla):
    textbox = slide.shapes.add_textbox(Inches(4.8), Inches(3.26), Inches(8), Inches(3))
    text_frame = textbox.text_frame

    p1 = text_frame.add_paragraph()
    p1_run = p1.add_run()
    p1_run.text = pantalla
    p1.font.bold = True
    p1.font.size = Pt(60)

def info_foto(slide, pantalla, elemento):
    textbox = slide.shapes.add_textbox(Inches(0.1299213), Inches(0.04), Inches(6), Inches(0.5))
    text_frame = textbox.text_frame

    p4 = text_frame.add_paragraph()
    p4.text = "CENTRO COMERCIAL: "
    p4.font.bold = True
    p4.font.size = Pt(16)
    p4.font.color.rgb = RGBColor(255, 0, 0)
    
    run2 = p4.add_run()
    run2.text = pantalla
    run2.font.bold = True
    run2.font.size = Pt(16)
    run2.font.color.rgb = RGBColor(0, 0, 0) 

    p5 = text_frame.add_paragraph()
    p5.text = "ELEMENTO: "
    p5.font.bold = True
    p5.font.size = Pt(16)
    p5.font.color.rgb = RGBColor(255, 0, 0)
    
    run3 = p5.add_run()
    run3.text = elemento
    run3.font.bold = True
    run3.font.size = Pt(16)
    run3.font.color.rgb = RGBColor(0, 0, 0) 

if __name__ == '__main__':
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    app.run(debug=True)
